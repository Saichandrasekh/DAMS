"""Generate a demo presentation deck for the Nandana DAMS pilot.
Produces both PPTX (editable) and PDF (ready-to-send) versions.

Run:
    .\\backend\\venv\\Scripts\\python.exe generate_demo_deck.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUT_DIR, 'DAMS_Nandana_Demo.pptx')
PDF_PATH = os.path.join(OUT_DIR, 'DAMS_Nandana_Demo.pdf')

# Brand colors (matching Nandana's purple primary)
PRIMARY = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1E, 0x29, 0x3B)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tx


def add_bullets(slide, left, top, width, height, items, size=16, color=DARK,
                line_spacing=1.3):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = '•  ' + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = 'Calibri'
    return tx


def add_rect(slide, left, top, width, height, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill(shp, color)
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    return shp


def add_header_bar(slide, title, subtitle=None):
    # Purple gradient bar at top
    bar = add_rect(slide, 0, 0, SW, Inches(0.9), PRIMARY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35),
                 subtitle, size=12, color=WHITE)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.25),
             'DAMS • Nandana School Pilot Demo', size=9, color=MUTED)
    add_text(slide, Inches(7), Inches(7.15), Inches(6), Inches(0.25),
             f'{page_no} / {total}', size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def feature_card(slide, left, top, width, height, icon_color, title, body):
    add_rect(slide, left, top, width, height, LIGHT_BG, line_color=MUTED)
    bar = add_rect(slide, left, top, Inches(0.12), height, icon_color)
    add_text(slide, left + Inches(0.3), top + Inches(0.15),
             width - Inches(0.4), Inches(0.45),
             title, size=15, bold=True, color=DARK)
    add_text(slide, left + Inches(0.3), top + Inches(0.6),
             width - Inches(0.4), height - Inches(0.7),
             body, size=11, color=DARK)


SLIDES = []


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────
def slide_title():
    s = add_blank()
    # full purple background
    add_rect(s, 0, 0, SW, SH, PRIMARY)
    # white card center
    add_rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5),
             WHITE, line_color=WHITE)
    add_text(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.6),
             'DAMS', size=42, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.5),
             'Digital Attendance Management System', size=22, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.4),
             'A complete school management platform', size=14,
             color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
             'Pilot proposal for', size=12, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.6),
             'Nandana School', size=28, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
             'Prepared 2026-05-20', size=10, color=WHITE,
             align=PP_ALIGN.CENTER)


SLIDES.append(slide_title)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────
def slide_problem():
    s = add_blank()
    add_header_bar(s, 'The problem', 'Why schools need a digital system')
    add_text(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
             'Paper registers and spreadsheets break down at scale',
             size=20, bold=True, color=DARK)
    items = [
        'Daily attendance lost in scattered registers — no way to spot patterns',
        'Parents only learn about absences at parent-teacher meetings, too late',
        'Fee dues tracked manually — receipts misplaced, defaulters missed',
        'No quick way for principal to see "who came today" across the school',
        'Teachers waste 15–20 minutes a day on register paperwork',
        'No central record of marks, holidays, announcements, or class diary',
    ]
    add_bullets(s, Inches(0.7), Inches(2.0), Inches(11.5), Inches(4.5),
                items, size=15)


SLIDES.append(slide_problem)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 3 — WHAT DAMS IS
# ─────────────────────────────────────────────────────────────────────────
def slide_what():
    s = add_blank()
    add_header_bar(s, 'What DAMS does',
                   'One platform — five logins — full school visibility')
    cards = [
        ('Principal / Admin', PRIMARY,
         'Full school dashboard. Add classes, students, teachers, fees. '
         'Real-time attendance and reports. Post school-wide announcements.'),
        ('Teacher', ACCENT,
         'Check-in / check-out, mark attendance per period, enter exam marks, '
         'post homework and class notes, view own attendance history.'),
        ('Student', SUCCESS,
         'Read attendance %, see timetable, view report card, check fee dues, '
         'read teacher homework and school announcements.'),
        ('Parent', WARNING,
         'See each child\'s attendance %, low-attendance alerts, recent activity, '
         'class diary entries and school announcements.'),
    ]
    card_w = Inches(2.85)
    card_h = Inches(3.5)
    for i, (title, color, body) in enumerate(cards):
        left = Inches(0.5 + i * 3.1)
        feature_card(s, left, Inches(1.6), card_w, card_h, color, title, body)
    add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5),
             'Built on modern web tech (React + Flask + SQLite). '
             'Runs on-premises or in the cloud. Works on any device — phone, '
             'tablet, or PC.',
             size=13, color=MUTED, align=PP_ALIGN.CENTER)


SLIDES.append(slide_what)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 4 — KEY MODULES
# ─────────────────────────────────────────────────────────────────────────
def slide_modules():
    s = add_blank()
    add_header_bar(s, 'Modules built',
                   'Everything a school needs, in one place')
    modules = [
        ('Attendance', 'Daily + per-period per-student tracking, with section drill-down and absentee filter'),
        ('Staff Attendance', 'Smart check-in/out for teachers, late detection, monthly summary'),
        ('Teachers Report', 'Per-teacher % attendance over any date range, with daily history'),
        ('Fees', 'Annual + monthly fee structure, payment collection, printable receipts, defaulters list'),
        ('Exams & Marks', 'Create exams, enter marks, publish report cards to students/parents'),
        ('Timetable', 'Per-class weekly schedule, teacher view, period-wise subject mapping'),
        ('Diary', 'School-wide announcements + class diary for homework and daily notes'),
        ('Batches', 'Year-cohort rollup view + Graduated alumni archive'),
        ('Classes & Promotion', 'Bulk year-end promotion, graduate flow, class history per student'),
    ]
    # 3x3 grid
    card_w = Inches(4.0)
    card_h = Inches(1.5)
    for i, (title, body) in enumerate(modules):
        row = i // 3
        col = i % 3
        left = Inches(0.4 + col * 4.2)
        top = Inches(1.5 + row * 1.7)
        feature_card(s, left, top, card_w, card_h, PRIMARY, title, body)


SLIDES.append(slide_modules)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 5 — NANDANA DEMO DATA (numbers)
# ─────────────────────────────────────────────────────────────────────────
def slide_demo_data():
    s = add_blank()
    add_header_bar(s, 'Nandana — demo dataset',
                   'A fully populated school ready for you to explore')
    stats = [
        ('10', 'Classes', 'Grades 6-10, sections A & B'),
        ('250', 'Students', '~25 per class with roll numbers'),
        ('15', 'Teachers', '3 per subject area'),
        ('50', 'Subjects', '5 per class'),
        ('420', 'Timetable slots', 'Mon-Sat × 7 periods'),
        ('21,000', 'Attendance records', 'Last 14 weekdays'),
        ('1,250', 'Exam marks', 'Mid-Term exam, published'),
        ('40', 'Fee heads', '4 per class (Tuition, Transport, etc.)'),
        ('781', 'Fee payments', 'Sample mix of paid/partial/pending'),
    ]
    # 3x3 stat tiles
    tile_w = Inches(4.0)
    tile_h = Inches(1.5)
    for i, (num, label, sub) in enumerate(stats):
        row = i // 3
        col = i % 3
        left = Inches(0.4 + col * 4.2)
        top = Inches(1.5 + row * 1.7)
        add_rect(s, left, top, tile_w, tile_h, LIGHT_BG, line_color=MUTED)
        add_text(s, left + Inches(0.2), top + Inches(0.15),
                 Inches(1.5), Inches(0.7),
                 num, size=28, bold=True, color=PRIMARY)
        add_text(s, left + Inches(1.8), top + Inches(0.2),
                 tile_w - Inches(2.0), Inches(0.45),
                 label, size=14, bold=True, color=DARK)
        add_text(s, left + Inches(1.8), top + Inches(0.65),
                 tile_w - Inches(2.0), tile_h - Inches(0.7),
                 sub, size=10, color=MUTED)


SLIDES.append(slide_demo_data)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 6 — PRINCIPAL DASHBOARD WALKTHROUGH
# ─────────────────────────────────────────────────────────────────────────
def slide_principal():
    s = add_blank()
    add_header_bar(s, 'Principal experience',
                   'What you see when you log in')
    items = [
        'Live attendance overview — present / absent / late counts for the whole school',
        'Click any class → see the period grid (P1-P7) for every student that day',
        'Staff tab in the same screen — mark teacher attendance, see who came / left',
        'Teachers Report — % attendance per teacher over any date range, drilldown to daily history',
        'Fees → see total demand, today\'s collection, top defaulters with phone numbers',
        'Batches → see the whole 2025-2026 cohort, plus Graduated alumni archive',
        'Diary → post school-wide announcements; see everything teachers post per class',
    ]
    add_bullets(s, Inches(0.7), Inches(1.4), Inches(11.5), Inches(5),
                items, size=14)
    add_text(s, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.4),
             'Login: principal@nandana.edu  /  nandana123',
             size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


SLIDES.append(slide_principal)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 7 — TEACHER EXPERIENCE
# ─────────────────────────────────────────────────────────────────────────
def slide_teacher():
    s = add_blank()
    add_header_bar(s, 'Teacher experience',
                   'Simple, fast, focused on the classroom')
    items = [
        'Dashboard with one-click Check-in / Check-out — late detection automatic',
        'My Classes — see assigned classes and subjects, with student counts',
        'Mark Attendance — pick class + subject + period + date, mark in seconds with bulk Set buttons',
        'Marks Entry — enter exam marks for any class/subject they teach',
        'Attendance Report — % per student in their class, identify low-attendance students',
        'Diary — post homework / daily notes for their classes; edit or delete their own posts',
        'My monthly attendance summary — present / late / absent / on-leave counts',
    ]
    add_bullets(s, Inches(0.7), Inches(1.4), Inches(11.5), Inches(5),
                items, size=14)
    add_text(s, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.4),
             'Sample login: gayatri.maths@nandana.edu  /  teacher123',
             size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


SLIDES.append(slide_teacher)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 8 — STUDENT / PARENT EXPERIENCE
# ─────────────────────────────────────────────────────────────────────────
def slide_student_parent():
    s = add_blank()
    add_header_bar(s, 'Student & Parent experience',
                   'Self-service portal — no more "what was the homework?"')
    # left half — Student
    add_text(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
             'Student', size=18, bold=True, color=SUCCESS)
    student_items = [
        'Overall attendance % with color-coded badge',
        'Subject-wise attendance breakdown',
        'Today\'s schedule with teacher names',
        'Monthly attendance calendar with filters',
        'Read-only weekly timetable',
        'Published exam marks (report card)',
        'My Fees: dues, payment history, print receipt',
        'Diary: teacher posts + school announcements',
    ]
    add_bullets(s, Inches(0.6), Inches(1.8), Inches(6.2), Inches(5),
                student_items, size=12)

    # right half — Parent
    add_text(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
             'Parent', size=18, bold=True, color=WARNING)
    parent_items = [
        'Each linked child shown as a card',
        'Live attendance % per child',
        'Low-attendance alerts (< 75%) flagged in red',
        'Recent attendance activity timeline',
        'View Details modal — full child report',
        'Diary: announcements + child\'s class entries',
        'Filter by which child (if multiple)',
        'No app install — works in any browser',
    ]
    add_bullets(s, Inches(6.9), Inches(1.8), Inches(6.2), Inches(5),
                parent_items, size=12)

    add_text(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
             'Student demo: n6a01@nandana.edu  /  student123       '
             '(Roll N6A01, Grade 6-A)',
             size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


SLIDES.append(slide_student_parent)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 9 — TECHNOLOGY & SECURITY
# ─────────────────────────────────────────────────────────────────────────
def slide_tech():
    s = add_blank()
    add_header_bar(s, 'Built on a proven stack',
                   'Reliable, secure, easy to maintain')
    # Two columns
    col_w = Inches(6)
    # left — tech
    add_text(s, Inches(0.5), Inches(1.3), col_w, Inches(0.4),
             'Technology', size=16, bold=True, color=PRIMARY)
    tech_items = [
        'Frontend: React 18 + TypeScript (mobile-friendly)',
        'Backend: Python Flask + Waitress',
        'Database: SQLite (zero setup) or PostgreSQL',
        'Authentication: JWT tokens + bcrypt hashing',
        'Web server: Nginx reverse proxy',
        'Deployment: Windows service or Linux container',
    ]
    add_bullets(s, Inches(0.6), Inches(1.8), col_w, Inches(3),
                tech_items, size=12)

    # right — security & operations
    add_text(s, Inches(6.8), Inches(1.3), col_w, Inches(0.4),
             'Security & operations', size=16, bold=True, color=PRIMARY)
    sec_items = [
        'Multi-tenant — each school\'s data fully isolated',
        'Rate-limited login (10/min) blocks brute force',
        'Daily automatic database backups (30-day retention)',
        'Audit log of every admin action',
        'Role-based access — teachers can\'t see other classes',
        'No plaintext passwords stored anywhere',
        'Auto-restart on crash via Windows services',
    ]
    add_bullets(s, Inches(6.9), Inches(1.8), col_w, Inches(3),
                sec_items, size=12)

    # bottom call-out
    add_rect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4),
             LIGHT_BG, line_color=PRIMARY)
    add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.4),
             'Already in production', size=14, bold=True, color=PRIMARY)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.8),
             'Live at attend.kautech.co.in. Same codebase, same features — '
             'just pointing at your school\'s data. Setup takes 1 day; full '
             'rollout (data import + training) takes 1 week.',
             size=12, color=DARK)


SLIDES.append(slide_tech)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 10 — DEMO LOGINS
# ─────────────────────────────────────────────────────────────────────────
def slide_logins():
    s = add_blank()
    add_header_bar(s, 'Demo credentials',
                   'Try the live demo with these accounts')
    rows = [
        ('Principal', 'principal@nandana.edu', 'nandana123', PRIMARY,
         'Full admin: classes, students, teachers, fees, diary, reports'),
        ('Teacher', 'gayatri.maths@nandana.edu', 'teacher123', ACCENT,
         'Mark attendance, enter marks, post diary, see own attendance'),
        ('Student', 'n6a01@nandana.edu', 'student123', SUCCESS,
         'Roll N6A01 — Grade 6-A. See attendance, marks, fees, diary'),
        ('Super Admin', 'superadmin@admin.com', 'admin123', WARNING,
         'Platform owner — manage multiple schools'),
    ]
    # header
    headers = ['Role', 'Email', 'Password', 'What they can do']
    col_widths = [Inches(2.0), Inches(3.8), Inches(2.0), Inches(5.0)]
    col_lefts = [Inches(0.4)]
    for w in col_widths[:-1]:
        col_lefts.append(col_lefts[-1] + w)

    top = Inches(1.4)
    row_h = Inches(0.5)
    add_rect(s, Inches(0.4), top, Inches(12.8), row_h, PRIMARY)
    for i, h in enumerate(headers):
        add_text(s, col_lefts[i] + Inches(0.1), top + Inches(0.1),
                 col_widths[i] - Inches(0.2), Inches(0.4),
                 h, size=13, bold=True, color=WHITE)

    # data rows
    for r_idx, (role, email, pwd, color, what) in enumerate(rows):
        row_top = top + row_h + Inches(0.05) + (Inches(0.95) * r_idx)
        bg = LIGHT_BG if r_idx % 2 == 0 else WHITE
        add_rect(s, Inches(0.4), row_top, Inches(12.8), Inches(0.9), bg)
        # left color stripe
        add_rect(s, Inches(0.4), row_top, Inches(0.12), Inches(0.9), color)
        add_text(s, col_lefts[0] + Inches(0.2), row_top + Inches(0.18),
                 col_widths[0] - Inches(0.3), Inches(0.5),
                 role, size=13, bold=True, color=color)
        add_text(s, col_lefts[1] + Inches(0.05), row_top + Inches(0.2),
                 col_widths[1] - Inches(0.1), Inches(0.5),
                 email, size=11, color=DARK)
        add_text(s, col_lefts[2] + Inches(0.05), row_top + Inches(0.2),
                 col_widths[2] - Inches(0.1), Inches(0.5),
                 pwd, size=11, bold=True, color=DARK)
        add_text(s, col_lefts[3] + Inches(0.05), row_top + Inches(0.15),
                 col_widths[3] - Inches(0.1), Inches(0.7),
                 what, size=10, color=MUTED)


SLIDES.append(slide_logins)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 11 — SUGGESTED DEMO FLOW
# ─────────────────────────────────────────────────────────────────────────
def slide_demo_flow():
    s = add_blank()
    add_header_bar(s, 'Suggested demo flow',
                   '5 minutes to show the full value')
    steps = [
        ('1', 'Log in as Principal — show the dashboard with school totals', PRIMARY),
        ('2', 'Open Attendance Today → click a class → switch to "Period grid" — show who came in which period', PRIMARY),
        ('3', 'Switch to Staff tab in same screen — mark a teacher present', PRIMARY),
        ('4', 'Open Teachers → Attendance Report tab → click any teacher → daily history modal', PRIMARY),
        ('5', 'Open Fees → pick a student → record a payment → print the receipt', ACCENT),
        ('6', 'Open Diary → post a school-wide announcement (e.g. "Sports Day")', ACCENT),
        ('7', 'Log out, log in as a Teacher (gayatri.maths@nandana.edu / teacher123)', SUCCESS),
        ('8', 'Click Check-in → see status badge update → post a homework entry in Diary', SUCCESS),
        ('9', 'Log out, log in as a Student (n6a01@nandana.edu / student123)', WARNING),
        ('10', 'Show their dashboard, fees page, diary with the announcement they just posted', WARNING),
    ]
    for i, (num, text, color) in enumerate(steps):
        row = i // 2
        col = i % 2
        left = Inches(0.4 + col * 6.3)
        top = Inches(1.4 + row * 1.05)
        add_rect(s, left, top, Inches(6.1), Inches(0.9), LIGHT_BG, line_color=MUTED)
        # number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1),
                                     top + Inches(0.2), Inches(0.5), Inches(0.5))
        fill(circle, color)
        ctf = circle.text_frame
        ctf.margin_left = 0; ctf.margin_right = 0
        ctf.margin_top = 0; ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = num
        crun.font.size = Pt(14)
        crun.font.bold = True
        crun.font.color.rgb = WHITE
        # text
        add_text(s, left + Inches(0.75), top + Inches(0.2),
                 Inches(5.2), Inches(0.6),
                 text, size=11, color=DARK)


SLIDES.append(slide_demo_flow)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 12 — NEXT STEPS
# ─────────────────────────────────────────────────────────────────────────
def slide_next():
    s = add_blank()
    add_header_bar(s, 'Next steps',
                   'How we move from demo to live deployment')
    steps = [
        ('Week 1', 'Pilot with one grade (e.g. Grade 8) — 50 students, 5 teachers, 1 admin. Run alongside paper register.'),
        ('Week 2', 'Import full student + teacher roster from your records (CSV upload supported).'),
        ('Week 3', 'Train teachers in 30-minute session. Train principal in 1-hour session.'),
        ('Week 4', 'Roll out school-wide. Paper register retired. SMS notifications optional.'),
        ('Ongoing', 'Daily DB backups, monthly review of system usage, feature requests prioritized.'),
    ]
    for i, (when, what) in enumerate(steps):
        top = Inches(1.5 + i * 0.95)
        add_rect(s, Inches(0.5), top, Inches(2.2), Inches(0.75),
                 PRIMARY)
        add_text(s, Inches(0.6), top + Inches(0.2),
                 Inches(2.0), Inches(0.4),
                 when, size=14, bold=True, color=WHITE)
        add_rect(s, Inches(2.7), top, Inches(10.1), Inches(0.75),
                 LIGHT_BG, line_color=MUTED)
        add_text(s, Inches(2.9), top + Inches(0.18),
                 Inches(9.7), Inches(0.5),
                 what, size=12, color=DARK)


SLIDES.append(slide_next)


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 13 — CONTACT / CALL TO ACTION
# ─────────────────────────────────────────────────────────────────────────
def slide_contact():
    s = add_blank()
    add_rect(s, 0, 0, SW, SH, PRIMARY)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.6),
             'Let\'s schedule a 30-minute walkthrough.', size=28, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(2.8), Inches(12.3), Inches(0.5),
             'See it on your own data. Bring your questions.',
             size=18, color=WHITE, align=PP_ALIGN.CENTER)
    # contact card
    add_rect(s, Inches(3.0), Inches(4.0), Inches(7.3), Inches(2.5),
             WHITE)
    add_text(s, Inches(3.0), Inches(4.3), Inches(7.3), Inches(0.4),
             'KAU Tech Services', size=18, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.0), Inches(4.8), Inches(7.3), Inches(0.4),
             'kautechservices@gmail.com', size=14, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.0), Inches(5.3), Inches(7.3), Inches(0.4),
             'Live demo: https://attend.kautech.co.in', size=14, color=DARK,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.0), Inches(5.8), Inches(7.3), Inches(0.4),
             'Source: github.com/Saichandrasekh/DAMS', size=12, color=MUTED,
             align=PP_ALIGN.CENTER)


SLIDES.append(slide_contact)


# Render
for fn in SLIDES:
    fn()

# Add footers (skip title and last contact slides)
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    if i in (1, total):
        continue
    add_footer(slide, i, total)

prs.save(PPTX_PATH)
print(f'OK  PPTX saved: {PPTX_PATH}')
print(f'    {total} slides')
print()
print('To convert to PDF, open the PPTX in PowerPoint and use File > Export > PDF.')
print('(Generating a true PDF requires PowerPoint or LibreOffice; the PPTX is the editable master.)')
